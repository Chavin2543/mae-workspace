#!/usr/bin/env python3
"""One-slide quarter performance summary in the Mitsui Fudosan (Asia)
"Thailand / Serviced Apartments" format (blue table + overview bullets).

Reads deck_data.json (produced by scripts/management_deck/extract_deck_data.py)
so every figure traces to the reconciled workbook / result sheets.

Usage:
  python3 scripts/management_deck/extract_deck_data.py /path/deck_data.json
  python3 scripts/quarter_summary.py /path/deck_data.json \
      --out "output/Thailand_SA_Q2-2026_summary.pptx" \
      --months 4-6 --label "Q2 2026 (Apr-Jun)"
"""
import argparse
import json

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

DAYS = [31, 28, 31, 30, 31, 30]
BLUE = RGBColor(0x5B, 0x9B, 0xD5)     # table header
ROW_A = RGBColor(0xDC, 0xE9, 0xF7)    # row fill
ROW_B = RGBColor(0xED, 0xF3, 0xFA)
INK = RGBColor(0x20, 0x20, 0x20)
NAVY = RGBColor(0x1E, 0x27, 0x61)
FLAG_RED = RGBColor(0xA5, 0x19, 0x31)
FLAG_NAVY = RGBColor(0x2D, 0x2A, 0x4A)

PROPS = [("SR9", "Somerset Rama 9"), ("AES", "Ascott Embassy Sathorn"),
         ("LYF", "Lyf Sukhumvit 8"), ("SP", "Somerset Pattaya")]
SEA = {"Indonesia", "Singapore", "Malaysia", "Vietnam", "Philippines"}
NAT_SHOW = {"China": "Chinese", "Thailand": "Thai"}


def is_long_stay(label):
    return "Long Stay" in label or label.strip().endswith("LS")


def metrics(d, m0, m1):
    """Per-property occ/ADR for months [m0..m1) + short-stay % + nat YTD."""
    out = {}
    for key, name in PROPS:
        p = d["perf"][key]
        ra = sold = rev = 0.0
        for i in range(m0, m1):
            ra_i = p["rooms"] * DAYS[i]
            ra += ra_i
            sold += p["occ"][i] * ra_i
            rev += p["occ"][i] * ra_i * p["adr"][i]
        seg = d["seg"][key]
        tot = sum(sum(v or 0 for v in r["rev"][m0:m1]) for r in seg)
        ls = sum(sum(v or 0 for v in r["rev"][m0:m1]) for r in seg if is_long_stay(r["label"]))
        tot_rn = sum(r["ytd"] for r in seg)
        nats = sorted(d["nat_rn"][key]["2026"], key=lambda e: -e["h1"])
        # group South-East Asia when it out-ranks single countries (LYF style)
        sea_share = sum(e["h1"] for e in nats if e["label"] in SEA) / tot_rn
        shown, used_sea = [], False
        for e in nats:
            lab = e["label"]
            if lab in SEA:
                if not used_sea and sea_share >= 0.15:
                    shown.append(("SEA", sea_share))
                    used_sea = True
                continue
            shown.append((NAT_SHOW.get(lab, lab), e["h1"] / tot_rn))
            if len(shown) >= 3:
                break
        out[key] = dict(name=name, occ=sold / ra, adr=rev / sold,
                        short=1 - ls / tot, nats=shown[:3])
    return out


def build(d, met, label, out_path):
    pres = Presentation()
    pres.slide_width, pres.slide_height = Inches(13.33), Inches(7.5)
    s = pres.slides.add_slide(pres.slide_layouts[6])

    def text(x, y, w, h, runs, align=PP_ALIGN.LEFT):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            for t, size, bold, color, underline in para:
                r = p.add_run()
                r.text = t
                r.font.size = Pt(size)
                r.font.bold = bold
                r.font.name = "Arial"
                r.font.color.rgb = color
                r.font.underline = underline
        return tb

    # header: Thailand + rule + Thai flag
    text(0.45, 0.15, 6, 0.75, [[("Thailand", 36, True, INK, False)]])
    ln = s.shapes.add_shape(1, Inches(0.45), Inches(0.95), Inches(12.45), Emu(19050))
    ln.fill.solid(); ln.fill.fore_color.rgb = INK; ln.line.fill.background()
    fx, fy, fw = 11.6, 0.15, 1.3
    for i, (col, hh) in enumerate([(FLAG_RED, 1), (RGBColor(255, 255, 255), 1),
                                   (FLAG_NAVY, 2), (RGBColor(255, 255, 255), 1),
                                   (FLAG_RED, 1)]):
        y0 = fy + sum([1, 1, 2, 1, 1][:i]) * 0.72 / 6
        r = s.shapes.add_shape(1, Inches(fx), Inches(y0), Inches(fw), Inches(hh * 0.72 / 6))
        r.fill.solid(); r.fill.fore_color.rgb = col
        r.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB); r.line.width = Pt(0.25)

    text(0.45, 1.05, 8, 0.4, [[("Serviced Apartments", 18, True, INK, True)]])
    text(0.45, 1.42, 8, 0.3, [[("Performance summary — " + label, 12, False,
                                RGBColor(0x60, 0x60, 0x60), False)]])

    # overview bullets
    A = d["arrivals"]
    mots = A["mots"]["ytd"]["2026"] / A["mots"]["ytd"]["2025"] - 1
    chn = A["chinese"]["ytd"]["2026"] / A["chinese"]["ytd"]["2025"] - 1
    me = A["mideast"]["ytd"]["2026"] / A["mideast"]["ytd"]["2025"] - 1
    pf = d["portfolio"]
    q2a, q2b = pf["act"][3:6], pf["bg"][3:6]
    occ = sum(m["sold"] for m in q2a) / sum(m["ra"] for m in q2a)
    occb = sum(m["sold"] for m in q2b) / sum(m["ra"] for m in q2b)
    adr = sum(m["revenue"] for m in q2a) / sum(m["sold"] for m in q2a)
    adrb = sum(m["revenue"] for m in q2b) / sum(m["sold"] for m in q2b)
    bullets = [
        [("Overview", 13, True, INK, True)],
        [("•  Tourist arrivals YTD below LY ({:+.1f}%, Jan–May); Q2 softened after the high season".format(mots * 100), 12, False, INK, False)],
        [("•  ", 12, False, INK, False), ("Middle East conflict", 12, True, INK, False),
         (" still the key disruptor ({:+.0f}% YTD arrivals) >> rate pressure".format(me * 100), 12, False, INK, False)],
        [("•  Chinese arrivals YTD {:+.0f}% vs last year — recovery continues".format(chn * 100), 12, False, INK, False)],
        [("•  Q2 volume & rate both under plan: portfolio OCC {:.1f}% (budget {:.1f}%), ADR THB {:,.0f} ({:+.1f}% vs budget); April was the soft spot (AES & Lyf, not the market)".format(
            occ * 100, occb * 100, adr, (adr / adrb - 1) * 100), 12, False, INK, False)],
    ]
    text(0.45, 1.75, 12.3, 1.6, bullets)

    # table
    headers = ["Property", "Occupancy", "Average Daily\nRate (ADR)", "Revenue\nContribution", "Occupants\nNationality YTD"]
    widths = [2.85, 1.55, 1.75, 1.95, 4.25]
    tbl_shape = s.shapes.add_table(5, 5, Inches(0.45), Inches(3.55), Inches(sum(widths)), Inches(3.3))
    tbl = tbl_shape.table
    for c, wdt in enumerate(widths):
        tbl.columns[c].width = Inches(wdt)
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255); p.font.name = "Arial"
    for r, (key, _) in enumerate(PROPS, start=1):
        m = met[key]
        nat1 = "{} {:.1f}%".format(*[m["nats"][0][0], m["nats"][0][1] * 100])
        nat2 = "  ".join("{} {:.1f}%".format(n, v * 100) for n, v in m["nats"][1:])
        cells = [m["name"], "{:.1f}%".format(m["occ"] * 100),
                 "THB {:,.0f}".format(m["adr"]),
                 "Short stay {:.0f}%".format(m["short"] * 100),
                 nat1 + "\n" + nat2]
        for c, val in enumerate(cells):
            cell = tbl.cell(r, c)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = ROW_A if r % 2 else ROW_B
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13); p.font.name = "Arial"; p.font.color.rgb = INK
                p.font.bold = (c == 0)

    text(7.6, 6.95, 5.3, 0.4, [[("MITSUI FUDOSAN (ASIA)", 20, True, INK, False)]],
         align=PP_ALIGN.RIGHT)
    pres.save(out_path)
    print("wrote", out_path)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("data", help="deck_data.json path")
    ap.add_argument("--out", required=True)
    ap.add_argument("--months", default="4-6", help="e.g. 4-6 for Q2")
    ap.add_argument("--label", default="Q2 2026 (Apr-Jun)")
    args = ap.parse_args()
    m0, m1 = [int(v) for v in args.months.split("-")]
    d = json.load(open(args.data))
    met = metrics(d, m0 - 1, m1)
    build(d, met, args.label, args.out)
    for key, _ in PROPS:
        m = met[key]
        print(key, "occ {:.1%} adr {:,.0f} short {:.0%}".format(m["occ"], m["adr"], m["short"]), m["nats"])


if __name__ == "__main__":
    main()
